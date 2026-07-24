# -*- coding: utf-8 -*-
"""금리 4팩터 — 팀원용 쉬운 설명 PPT (make_rates_ppt.py 의 입문 버전).

정식 운용 설명서(make_rates_ppt.py)가 "무엇을 어떻게 하는가"를 빠짐없이 적은
문서라면, 이 덱은 "처음 보는 사람이 이해할 수 있는가"만 목표로 한다.
- 팩터 4개를 비유로 먼저 설명하고 수식은 최소화
- 오늘 실제 숫자로 포지션이 정해지는 과정을 STEP 1~5 로 따라간다

⚠ STEP 예시의 모든 숫자는 2026-07-23 종가 기준 엔진 실측값이다 (임의로 만든
   값이 아님). 숫자를 갱신하려면 아래 EX_* 상수를 sleeve_engine 산출로 다시
   뽑아 넣을 것 — 중간 단계는 _combine_class / _realized_vol /
   compute_target_positions / finalize_positions 에서 나온다.

Usage: python scripts/make_rates_ppt_easy.py [-o 출력경로.pptx]
"""
import argparse
from datetime import date

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from make_rates_ppt import (          # 디자인 시스템 재사용 (미래에셋 CI)
    NAVY, ORANGE, WHITE, INK, MUTED, CARD_BG, CARD_LN, RED, GREEN, FONT, W, H,
    _run, _para, _rect, _textbox, add_slide, card, table, note)

# ── 2026-07-23 종가 기준 실측 체인 ──────────────────────────────────────
ASSETS = ['미국 2Y', '미국 10Y', '한국 3Y', '한국 10Y']
EX_FACTORS = {                        # 슬리브 스냅샷 (중립화·게이트 적용 후)
    '추세':      [-2.19, -1.96, -1.14, -1.25],
    '밸류':      [-0.26, +0.07, +0.34, +0.51],
    '캐리':      [-0.02, -0.00, -0.02, +0.25],
    '정책':      [-1.96, -1.96, -0.91, -0.91],
}
EX_CONV   = [-1.028, -0.885, -0.345, -0.260]   # 4개 평균 → EWMA 5일
EX_VOL    = ['1.3%', '4.3%', '2.4%', '6.7%']
EX_INV    = [7.70, 2.33, 4.09, 1.49]
EX_MULT   = [-7.907, -2.066, -1.412, -0.389]
EX_VT     = [-1.403, -0.966, -0.660, -0.182]   # 볼타겟 스케일 0.1775
EX_SMOOTH = [-1.530, -0.996, -0.562, -0.135]
EX_FINAL  = [-0.765, -0.498, -0.281, -0.068]   # × 노출스케일 0.5
EX_PREV   = [-0.781, -0.502, -0.269, -0.062]
PER_UNIT  = 1252.0                              # 포지션 1.0 = 1,252만원


def _f(vals, fmt='+.2f'):
    return [format(v, fmt) if isinstance(v, (int, float)) else v for v in vals]


def step_badge(slide, x, y, n, label, w=Inches(2.6)):
    """STEP n 배지 + 라벨."""
    _rect(slide, x, y, w, Inches(0.42), fill=NAVY)
    tf = _textbox(slide, x + Inches(0.10), y + Inches(0.02), w - Inches(0.2),
                  Inches(0.38), anchor=MSO_ANCHOR.MIDDLE)
    p = _para(tf, True)
    _run(p, f'STEP {n}', size=11, bold=True, color=ORANGE)
    _run(p, f'   {label}', size=12, bold=True, color=WHITE)


def big(slide, x, y, w, h, value, caption, color=NAVY, vsize=30):
    """큰 숫자 강조 박스."""
    _rect(slide, x, y, w, h, fill=CARD_BG, line=CARD_LN)
    tf = _textbox(slide, x, y + Inches(0.12), w, h - Inches(0.2),
                  anchor=MSO_ANCHOR.MIDDLE)
    p = _para(tf, True, align=PP_ALIGN.CENTER)
    _run(p, value, size=vsize, bold=True, color=color)
    p = _para(tf, False, space_before=4, align=PP_ALIGN.CENTER)
    _run(p, caption, size=10, color=MUTED)


def arrow(slide, x, y, size=Inches(0.34)):
    a = _rect(slide, x, y, size, size, fill=ORANGE, shape=MSO_SHAPE.RIGHT_ARROW)
    return a


# ═══════════════════════════════════════════════════════════════════════
def s1_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    _rect(s, Inches(0.9), Inches(2.45), Inches(1.5), Inches(0.075), fill=ORANGE)

    tf = _textbox(s, Inches(0.9), Inches(2.65), Inches(11.4), Inches(1.6))
    p = _para(tf, True)
    _run(p, '금리 4팩터, 쉽게 이해하기', size=38, bold=True, color=WHITE)
    p = _para(tf, False, space_before=6)
    _run(p, '오늘 포지션이 정해지는 과정을 숫자로 따라가 봅니다', size=20, color=ORANGE)

    tf2 = _textbox(s, Inches(0.92), Inches(4.6), Inches(11.0), Inches(1.3))
    p = _para(tf2, True, line=1.4)
    _run(p, '수식 없이 4개 팩터가 무엇을 보는지 이해하고,',
         size=14, color=RGBColor(0xBF, 0xD4, 0xEA))
    p = _para(tf2, False, space_before=6, line=1.4)
    _run(p, '실제 하루치 데이터로 미국·한국 국채선물 포지션이 어떻게 나오는지 '
            '단계별로 확인합니다.', size=14, color=RGBColor(0xBF, 0xD4, 0xEA))

    _rect(s, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.02),
          fill=RGBColor(0x2A, 0x5A, 0x8C))
    tf3 = _textbox(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.5))
    p = _para(tf3, True)
    _run(p, f'{date.today().isoformat()}   |   예시 데이터 2026-07-23 종가 기준 '
            f'(실제 엔진 산출값)   |   내부 교육 자료',
         size=10, color=RGBColor(0x8F, 0xAD, 0xCC))


def s2_signconv(prs):
    s = add_slide(prs, '시작 전에 — 딱 하나만 기억하세요',
                  '이걸 헷갈리면 이후 모든 숫자가 거꾸로 읽힙니다', 2)
    big(s, Inches(0.45), Inches(1.35), Inches(5.9), Inches(1.75),
        '채권선물 롱 (+)  =  금리 하락 베팅', '', NAVY, 22)
    big(s, Inches(6.95), Inches(1.35), Inches(5.9), Inches(1.75),
        '채권선물 숏 (−)  =  금리 상승 베팅', '', RED, 22)

    card(s, Inches(0.45), Inches(3.35), Inches(12.4), Inches(1.55),
         '왜 그런가요?',
         ['채권 가격과 금리는 항상 반대로 움직입니다. 금리가 내리면 기존에 발행된 '
          '(이자를 더 많이 주는) 채권의 가치가 올라가기 때문입니다.',
          ('그래서 "국채선물을 산다(롱)" = "앞으로 금리가 내릴 것 같다"는 뜻입니다. '
           '이 자료의 모든 + / − 부호는 이 규칙을 따릅니다.', NAVY, True)])

    card(s, Inches(0.45), Inches(5.10), Inches(12.4), Inches(1.45),
         '오늘의 예를 미리 보면',
         [('오늘 이 북의 포지션은 4종목 전부 마이너스(−)입니다 → '
           '"한국과 미국 금리가 앞으로 오를 것 같다"에 베팅하고 있다는 뜻입니다.',
           ORANGE, True)])


def s3_onesentence(prs):
    s = add_slide(prs, '이 북을 한 문장으로',
                  '서로 다른 4명에게 물어보고, 의견을 모아서, 크기를 정한다', 3)
    tf = _textbox(s, Inches(0.45), Inches(1.30), Inches(12.4), Inches(0.9))
    p = _para(tf, True, align=PP_ALIGN.CENTER, line=1.3)
    _run(p, '"금리가 오를까 내릴까?" 를 4가지 다른 방식으로 묻고,\n'
            '답을 평균 내어 한국·미국 국채선물의 매수/매도 크기를 정한다.',
         size=17, bold=True, color=NAVY)

    names = [('추세', '요즘 어느 쪽으로\n움직이고 있나?'),
             ('밸류', '평소 가격에 비해\n비싼가 싼가?'),
             ('캐리', '들고 있으면 이자를\n얼마나 주나?'),
             ('정책', '중앙은행은 어디로\n가고 있나?')]
    x = Inches(0.45)
    for i, (nm, q) in enumerate(names):
        cx = x + Inches(3.19) * i
        _rect(s, cx, Inches(2.45), Inches(2.95), Inches(1.85),
              fill=CARD_BG, line=CARD_LN)
        _rect(s, cx, Inches(2.45), Inches(2.95), Inches(0.05), fill=ORANGE)
        tf2 = _textbox(s, cx, Inches(2.60), Inches(2.95), Inches(1.6),
                       anchor=MSO_ANCHOR.MIDDLE)
        p = _para(tf2, True, align=PP_ALIGN.CENTER)
        _run(p, nm, size=19, bold=True, color=NAVY)
        p = _para(tf2, False, space_before=8, align=PP_ALIGN.CENTER, line=1.25)
        _run(p, q, size=11, color=INK)

    card(s, Inches(0.45), Inches(4.55), Inches(12.4), Inches(1.95),
         '중요한 것 두 가지',
         [('① 4명의 의견은 똑같은 비중(각 25%)으로 반영됩니다. '
           '어느 한 명을 더 믿지 않습니다.', NAVY, True),
          ('② 의견의 세기는 "방향"만 정합니다. 크기는 별도의 리스크 계산이 정합니다 '
           '— 신호가 2배 강해져도 포지션이 2배가 되지 않습니다.', NAVY, True),
          '이 두 가지가 이 북이 "예측을 잘해서"가 아니라 "규칙을 지켜서" '
          '돈을 버는 구조인 이유입니다.'])


def _factor_slide(prs, page, num, title, oneline, detail, when_long, kind,
                  kind_desc, why):
    s = add_slide(prs, f'팩터 {num} — {title}', oneline, page)
    card(s, Inches(0.45), Inches(1.32), Inches(7.9), Inches(2.85),
         '쉽게 말하면', detail, body_size=12)

    # 성격 카드 — 큰 라벨 + 설명을 한 박스에 담는다 (라벨만 떠 있으면 미아처럼 보임)
    acc = NAVY if kind == '방향성' else ORANGE
    _rect(s, Inches(8.75), Inches(1.32), Inches(4.1), Inches(2.85),
          fill=CARD_BG, line=CARD_LN)
    _rect(s, Inches(8.75), Inches(1.32), Inches(0.045), Inches(2.85), fill=acc)
    tf = _textbox(s, Inches(8.95), Inches(1.45), Inches(3.75), Inches(2.6))
    p = _para(tf, True, align=PP_ALIGN.CENTER)
    _run(p, '이 팩터의 성격', size=11, color=MUTED)
    p = _para(tf, False, space_before=8, align=PP_ALIGN.CENTER)
    _run(p, kind, size=26, bold=True, color=acc)
    p = _para(tf, False, space_before=14, line=1.35)
    _run(p, kind_desc, size=11, color=INK)

    card(s, Inches(0.45), Inches(4.40), Inches(6.1), Inches(1.85),
         '언제 채권을 사나 (롱)', when_long, body_size=11.5)
    card(s, Inches(6.75), Inches(4.40), Inches(6.1), Inches(1.85),
         '왜 돈이 되나', why, body_size=11.5)
    return s


def s4_trend(prs):
    s = _factor_slide(
        prs, 4, '①', '추세 (Trend)', '요즘 어느 쪽으로 움직이고 있나?',
        ['최근 6개월과 1년 동안 채권 가격이 오르는 흐름이었다면 '
         '"그 흐름이 좀 더 갈 것"에 베팅합니다.',
         ('달리던 방향으로 계속 달린다 — 가장 고전적인 추세추종입니다.',
          NAVY, True),
         ('단, 방향 없이 오르락내리락하는 장세로 판단되면 이 팩터를 '
          '자동으로 꺼버립니다 (레짐 게이트).', ORANGE, True)],
        ['최근 6~12개월 채권 가격이 오르는 추세일 때',
         '= 금리가 꾸준히 내려오고 있을 때'],
        '방향성',
        '다른 나라와 비교하지 않고, 그 나라 자체의 흐름만 봅니다. '
        '8개국이 다 같이 오르면 다 같이 삽니다.',
        ['금리 사이클은 몇 달 단위로 관성이 있습니다. 중앙은행 정책도, '
         '경기 흐름도 하루아침에 방향을 바꾸지 않기 때문입니다.'])
    note(s, Inches(6.40),
         '※ 너무 짧은 추세(3개월 이하)는 방향이 자주 뒤집혀 손실만 쌓였습니다. 그래서 6개월·1년만 씁니다.')


def s5_value(prs):
    s = _factor_slide(
        prs, 5, '②', '밸류 (Value)', '평소 가격에 비해 비싼가 싼가?',
        ['지난 2년 평균 가격과 비교해서 지금이 비싼지 싼지를 봅니다. '
         '너무 오른 채권은 팔고, 너무 내린 채권은 삽니다.',
         ('고무줄처럼 — 많이 늘어나면 언젠가 제자리로 돌아온다는 생각입니다.',
          NAVY, True),
         ('단, "절대적으로 싸다"가 아니라 "8개국 중에서 상대적으로 싸다"를 봅니다.',
          ORANGE, True)],
        ['2년 평균보다 싼 편일 때',
         '(= 다른 나라들보다 금리가 더 많이 올라온 상태)'],
        '상대가치',
        '8개국 평균을 빼고 비교합니다. 모든 나라가 다 같이 싸면 아무도 사지 '
        '않습니다 — 오직 순위만 봅니다.',
        ['시장은 한쪽으로 과하게 쏠렸다가 되돌아오는 경향이 있습니다. '
         '추세 팩터가 놓치는 "너무 갔다" 구간을 이 팩터가 잡아줍니다.'])
    note(s, Inches(6.40),
         '※ 추세와 밸류는 서로 반대 방향을 가리킬 때가 많습니다. 일부러 그렇게 섞었습니다 — 한쪽이 틀릴 때 다른 쪽이 버텨줍니다.')


def s6_carry(prs):
    s = _factor_slide(
        prs, 6, '③', '캐리 (Carry)', '들고 있으면 이자를 얼마나 주나?',
        ['가격이 전혀 안 움직여도, 채권을 들고 있으면 이자가 나옵니다. '
         '그 이자가 큰 나라의 채권을 삽니다.',
         ('단, 환헤지 비용을 뺀 "실제로 남는 이자"로 계산합니다. '
          '금리가 높아 보여도 환헤지 비용이 다 먹으면 남는 게 없습니다.',
          ORANGE, True),
         ('예금 이자율이 높은 은행에 돈을 넣되, 수수료를 뺀 실수령액으로 '
          '비교하는 것과 같습니다.', NAVY, True)],
        ['환헤지 비용을 빼고도 이자가 많이 남는 나라일 때'],
        '상대가치',
        '8개국 평균을 빼고 비교합니다. "이자를 많이 주는 나라 롱 / '
        '적게 주는 나라 숏"만 남습니다.',
        ['시장이 아무 방향으로도 안 움직여도 이자만큼은 쌓입니다. '
         '가장 꾸준하고 설명하기 쉬운 수익원입니다.'])
    note(s, Inches(6.40),
         '※ 참고: 미국 2Y와 한국 3Y처럼 만기가 짧은 종목은 계산 구조상 캐리 값이 똑같이 나옵니다 — 이 팩터는 그 둘을 구분하지 못합니다.')


def s7_policy(prs):
    s = _factor_slide(
        prs, 7, '④', '정책 모멘텀 (Policy)', '중앙은행은 어디로 가고 있나?',
        ['그 나라의 짧은 만기 금리(2~3년)가 지난 6개월간 오르고 있었는지 '
         '내리고 있었는지를 봅니다.',
         ('짧은 금리는 중앙은행 정책을 가장 빠르게 반영합니다 — '
          '내려오고 있다 = 금리인하 사이클 = 채권 매수.', NAVY, True),
         ('가격이 아니라 "금리" 자체를 직접 보는 유일한 팩터입니다.',
          ORANGE, True)],
        ['단기금리가 6개월간 내려왔을 때',
         '(= 중앙은행이 완화 쪽으로 가고 있을 때)'],
        '방향성',
        '나라별로만 구분하고 만기는 구분하지 않습니다 — 미국 2Y와 미국 10Y는 '
        '항상 같은 값을 받습니다.',
        ['중앙은행 사이클은 한번 방향을 잡으면 1~2년씩 갑니다. '
         '전략 공장에서도 가장 성과가 좋았던 아이디어입니다.'])
    note(s, Inches(6.40),
         '※ 오늘 미국 정책 값은 −1.96 입니다 — 미국 단기금리가 6개월간 꽤 올랐다(긴축)는 뜻이고, 그래서 강한 매도 신호를 냅니다.')


def s8_compare(prs):
    s = add_slide(prs, '4팩터 한눈에 비교', '무엇이 다르고, 왜 4개를 같이 쓰나', 8)
    table(s, Inches(0.45), Inches(1.30), Inches(12.4),
          ['팩터', '무엇을 보나', '언제 사나 (롱)', '성격', '기간'],
          [['① 추세', '최근 6개월·1년 가격 흐름', '오르는 추세일 때', '방향성', '6~12개월'],
           ['② 밸류', '2년 평균 대비 비싼가 싼가', '남들보다 싼 편일 때', '상대가치', '2년'],
           ['③ 캐리', '환헤지 후 실제 남는 이자', '이자가 많이 남을 때', '상대가치', '현재 수준'],
           ['④ 정책', '단기금리의 6개월 변화', '중앙은행이 완화할 때', '방향성', '6개월']],
          col_w=[1.0, 3.0, 2.4, 1.1, 1.2], fs=11,
          row_h=Inches(0.52), head_h=Inches(0.42))

    card(s, Inches(0.45), Inches(3.90), Inches(6.1), Inches(1.35),
         '"방향성" 팩터 — 추세 · 정책',
         [('금리가 오를까 내릴까에 직접 베팅합니다. 8개국이 다 같이 사라고 하면 '
           '다 같이 삽니다.', NAVY, True)])
    card(s, Inches(6.75), Inches(3.90), Inches(6.1), Inches(1.35),
         '"상대가치" 팩터 — 밸류 · 캐리',
         [('나라끼리 순위만 매깁니다. 8개국 평균을 빼기 때문에 '
           '"어디가 남들보다 나은가"만 남습니다.', ORANGE, True)])

    card(s, Inches(0.45), Inches(5.42), Inches(12.4), Inches(1.15),
         '왜 4개를 같이 쓰나',
         [('네 팩터는 서로 다른 것을 보기 때문에 동시에 틀릴 확률이 낮습니다. '
           '한두 개가 헛다리를 짚어도 나머지가 버텨줍니다 — 이것이 분산의 핵심입니다.',
           NAVY, True)])


def s9_jargon(prs):
    s = add_slide(prs, '숫자 읽는 법 — 딱 두 가지 개념',
                  'z-score 와 "평균 빼기(중립화)"', 9)
    card(s, Inches(0.45), Inches(1.32), Inches(6.1), Inches(2.60),
         'z-score = "평소보다 얼마나 특이한가"',
         ['모든 팩터 값은 z-score로 바뀝니다. 단위가 다른 것들(가격·이자·금리변화)을 '
          '같은 자로 재기 위해서입니다.',
          ('0  =  평소와 똑같음', NAVY, True),
          ('+1 =  평소보다 꽤 높은 편  /  +2 = 매우 드물게 높음', NAVY, True),
          ('−2 =  매우 드물게 낮음  (오늘 미국 2Y 추세가 −2.19)', RED, True),
          '기준은 "지난 1년 동안의 자기 자신"입니다.'])
    card(s, Inches(6.75), Inches(1.32), Inches(6.1), Inches(2.60),
         '중립화 = "8개국 평균을 뺀다"',
         ['밸류·캐리에만 적용합니다. 평균을 빼면 "다 같이 오를까"는 사라지고 '
          '"누가 남들보다 나은가"만 남습니다.',
          ('예: 오늘 미국 2Y의 밸류 원값은 +1.73 (싼 편)이지만,', INK, False),
          ('8개국이 더 싸서 평균을 빼면 −0.26 이 됩니다 → 상대적으로는 '
           '비싼 편.', ORANGE, True),
          '추세·정책은 평균을 빼지 않습니다 — 방향 베팅을 그대로 남기기 위해서.'])

    card(s, Inches(0.45), Inches(4.10), Inches(12.4), Inches(2.10),
         '⚠ 여기서 자주 나오는 오해',
         [('"평균을 뺐으니 서로 상쇄돼서 위험이 없는 것 아닌가요?"', NAVY, True),
          '아닙니다. 평균은 8개국(영국·일본·호주 포함)으로 계산하는데, 실제로 '
          '주문하는 건 한국·미국 4종뿐입니다.',
          ('즉 상대가치의 "반대편 다리"를 실제로는 사지 않기 때문에, 일부 방향 '
           '노출이 남습니다. 의도된 설계이고, 검증에서도 이 편이 더 나았습니다.',
           ORANGE, True)])


def s10_step1(prs):
    s = add_slide(prs, '예시 STEP 1 — 오늘 4명은 뭐라고 했나',
                  '2026-07-23 종가 기준 · 실제 엔진 값', 10)
    step_badge(s, Inches(0.45), Inches(1.26), 1, '4팩터 의견 수집')

    rows = [[k] + _f(v) for k, v in EX_FACTORS.items()]
    table(s, Inches(0.45), Inches(1.86), Inches(8.6),
          ['팩터'] + ASSETS, rows,
          col_w=[1.3, 1.15, 1.15, 1.15, 1.15], fs=12,
          row_h=Inches(0.52), head_h=Inches(0.44))

    card(s, Inches(9.35), Inches(1.86), Inches(3.5), Inches(2.52),
         '읽어보면',
         [('추세·정책은 4종목 전부 마이너스', RED, True),
          '= 둘 다 "금리 오른다"에 강하게 베팅',
          ('밸류·캐리는 한국 쪽이 플러스', GREEN, True),
          '= 한국이 상대적으로 매력적'])

    card(s, Inches(0.45), Inches(4.60), Inches(12.4), Inches(1.75),
         '이 표에서 바로 보이는 것',
         [('미국 2Y의 추세 −2.19 와 정책 −1.96 이 오늘 가장 강한 신호입니다 — '
           '미국 단기금리가 6개월간 뚜렷하게 올랐고, 가격도 계속 내려왔다는 뜻.',
           NAVY, True),
          '반대로 한국 10Y는 밸류 +0.51, 캐리 +0.25 로 "사라"는 의견이지만 '
          '크기가 작습니다.',
          ('→ 결과적으로 4종목 모두 매도 방향이 되고, 그중 미국이 더 강하게 나올 '
           '것임을 여기서 이미 짐작할 수 있습니다.', ORANGE, True)])


def s11_step2(prs):
    s = add_slide(prs, '예시 STEP 2 — 의견을 하나로 모은다',
                  '4명의 평균 = 컨빅션 (얼마나 확신하는가)', 11)
    step_badge(s, Inches(0.45), Inches(1.26), 2, '단순 평균 → 5일 부드럽게',
               w=Inches(3.3))          # 라벨이 길어 기본 폭(2.6")을 넘김

    tf = _textbox(s, Inches(0.45), Inches(1.88), Inches(12.4), Inches(0.55))
    p = _para(tf, True, align=PP_ALIGN.CENTER)
    _run(p, '미국 2Y :  ( −2.19  +  (−0.26)  +  (−0.02)  +  (−1.96) )  ÷  4  '
            '=  −1.11', size=15, bold=True, color=NAVY)

    table(s, Inches(0.45), Inches(2.62), Inches(12.4),
          ['단계'] + ASSETS,
          [['4팩터 단순 평균'] + _f([-1.11, -0.96, -0.43, -0.35]),
           ['최근 5일 평활 후 (컨빅션)'] + _f(EX_CONV, '+.3f')],
          col_w=[2.6, 1.3, 1.3, 1.3, 1.3], fs=12,
          row_h=Inches(0.52), head_h=Inches(0.44),
          cell_colors={(1, c): NAVY for c in range(5)})

    card(s, Inches(0.45), Inches(4.32), Inches(6.1), Inches(1.65),
         '왜 5일 평활을 하나',
         ['하루치 신호가 튀었다고 바로 반응하면 매매만 늘고 비용만 나갑니다.',
          ('최근 5일을 부드럽게 섞어 "진짜 바뀐 것"만 남깁니다.', NAVY, True)])
    card(s, Inches(6.75), Inches(4.32), Inches(6.1), Inches(1.65),
         '컨빅션이 뜻하는 것',
         [('−1.03 은 "미국 2Y를 팔자"는 방향과 그 확신의 세기입니다.', NAVY, True),
          '아직 "몇 계약 팔지"는 정해지지 않았습니다 — 그건 STEP 3부터.'])

    note(s, Inches(6.15),
         '※ 4개를 단순 평균하는 이유: 어느 팩터가 앞으로 더 잘 맞을지 미리 알 수 없기 때문입니다. '
         '과거 성과를 보고 가중치를 조정하면 과최적화가 되어 실전에서 무너집니다.')


def s12_step3(prs):
    s = add_slide(prs, '예시 STEP 3 — 크기 정하기 ①  "덜 흔들리면 더 많이"',
                  '같은 확신이라도 종목마다 계약 수가 다른 이유', 12)
    step_badge(s, Inches(0.45), Inches(1.26), 3, '변동성으로 나누기')

    table(s, Inches(0.45), Inches(1.86), Inches(12.4),
          ['단계'] + ASSETS,
          [['컨빅션 (STEP 2)'] + _f(EX_CONV, '+.3f'),
           ['연 변동성 (얼마나 흔들리나)'] + EX_VOL,
           ['→ 배수 (덜 흔들릴수록 크게)'] + _f(EX_INV, '.2f'),
           ['컨빅션 × 배수'] + _f(EX_MULT, '+.2f')],
          col_w=[3.0, 1.3, 1.3, 1.3, 1.3], fs=12,
          row_h=Inches(0.50), head_h=Inches(0.44),
          cell_colors={(2, c): ORANGE for c in range(1, 5)})

    card(s, Inches(0.45), Inches(4.52), Inches(7.9), Inches(1.95),
         '여기가 가장 헷갈리는 부분입니다',
         [('미국 2Y는 하루에 잘 안 움직입니다 (연 변동성 1.3%). '
           '그래서 같은 확신이라도 7.70배 많이 삽니다.', NAVY, True),
          ('한국 10Y는 많이 움직입니다 (6.7%). 그래서 1.49배만 삽니다.',
           NAVY, True),
          ('→ 명목 금액이 크다고 위험이 큰 게 아닙니다. '
           '모든 종목이 비슷한 크기의 위험을 지도록 맞춘 결과입니다.',
           ORANGE, True)])

    big(s, Inches(8.75), Inches(4.52), Inches(4.1), Inches(1.95),
        '위험 균등', '명목이 아니라 "흔들리는 정도"를 기준으로\n'
        '모든 종목에 같은 위험을 배분합니다', NAVY, 26)


def s13_step4(prs):
    s = add_slide(prs, '예시 STEP 4 — 크기 정하기 ②  마지막 3번의 조정',
                  '북 전체 위험 맞추기 → 천천히 움직이기 → 절반으로 줄이기', 13)
    step_badge(s, Inches(0.45), Inches(1.26), 4, '전체 조정')

    table(s, Inches(0.45), Inches(1.86), Inches(12.4),
          ['조정'] + ASSETS + ['무엇을 하는가'],
          [['컨빅션 × 배수'] + _f(EX_MULT, '+.2f') + ['STEP 3 결과'],
           ['① 북 전체 위험 맞추기'] + _f(EX_VT, '+.3f') + ['전체를 약 18%로 축소'],
           ['② 천천히 움직이기'] + _f(EX_SMOOTH, '+.3f') + ['어제 포지션 80% 유지'],
           ['③ 최종 절반으로'] + _f(EX_FINAL, '+.3f') + ['운용 규모에 맞춤']],
          col_w=[2.3, 1.1, 1.1, 1.1, 1.1, 2.6], fs=11.5,
          row_h=Inches(0.50), head_h=Inches(0.44),
          cell_colors={(3, c): NAVY for c in range(5)})

    card(s, Inches(0.45), Inches(4.52), Inches(4.0), Inches(1.95),
         '① 북 전체 위험 맞추기',
         [('금리 북 전체가 연 7.1% 정도만 흔들리도록 전체 크기를 조절합니다.',
           NAVY, True),
          '시장이 요동치면 자동으로 작아지고, 잠잠하면 커집니다.'])
    card(s, Inches(4.65), Inches(4.52), Inches(4.0), Inches(1.95),
         '② 천천히 움직이기',
         [('오늘 새 목표의 20%만 반영하고 80%는 어제 것을 유지합니다.',
           NAVY, True),
          '신호가 출렁여도 실제 주문은 잘게 나옵니다 — 비용 절감의 핵심.'])
    card(s, Inches(8.85), Inches(4.52), Inches(4.0), Inches(1.95),
         '③ 손실이 커지면?',
         [('북 손실이 4%를 넘으면 자동으로 절반, 8%를 넘으면 전부 청산합니다.',
           RED, True),
          '오늘은 손실이 없어 그대로 통과했습니다.'], accent=RED)


def s14_final(prs):
    s = add_slide(prs, '예시 STEP 5 — 오늘의 최종 주문',
                  '여기까지가 매일 아침 대시보드에 뜨는 숫자입니다', 14)
    step_badge(s, Inches(0.45), Inches(1.26), 5, '주문 확정')

    deltas = [f'{v * PER_UNIT:+,.0f}' for v in EX_FINAL]
    diffs = [f'{(a - b):+.3f}' for a, b in zip(EX_FINAL, EX_PREV)]
    table(s, Inches(0.45), Inches(1.86), Inches(12.4),
          ['항목'] + ASSETS + ['합계'],
          [['오늘 목표 포지션'] + _f(EX_FINAL, '+.3f') + [f'{sum(EX_FINAL):+.2f}'],
           ['어제 포지션'] + _f(EX_PREV, '+.3f') + [f'{sum(EX_PREV):+.2f}'],
           ['오늘 주문 (차이)'] + diffs + ['—'],
           ['금액 환산 (만원)'] + deltas +
            [f'{sum(EX_FINAL) * PER_UNIT:+,.0f}']],
          col_w=[2.2, 1.25, 1.25, 1.25, 1.25, 1.3], fs=11.5,
          row_h=Inches(0.50), head_h=Inches(0.44),
          cell_colors={(0, c): NAVY for c in range(6)})

    card(s, Inches(0.45), Inches(4.52), Inches(7.9), Inches(1.95),
         '오늘 이 포지션인 이유 — 한 문단으로',
         [('한국·미국 국채선물 4종을 모두 매도(금리 상승 베팅)하고 있습니다.',
           NAVY, True),
          '추세와 정책 두 팩터가 강하게 매도를 가리켰고(미국 −2.19 / −1.96), '
          '밸류·캐리가 한국 쪽에 소폭 매수를 냈지만 이기지 못했습니다.',
          ('미국 매도가 더 큰 이유는 신호가 더 강했고(정책 −1.96 vs −0.91), '
           '미국 2Y가 덜 흔들려 같은 위험에 더 많은 계약이 들어가기 때문입니다.',
           ORANGE, True)])

    card(s, Inches(8.75), Inches(4.52), Inches(4.1), Inches(1.95),
         '주문은 이만큼뿐',
         [('어제와 비교하면 변화는 미미합니다 (최대 +0.02).', NAVY, True),
          '"천천히 움직이기" 덕분에 매일 큰 거래를 하지 않습니다.',
          ('연간 회전율 17배 — 하루 평균으로는 아주 작은 주문입니다.',
           ORANGE, True)])


def s15_faq(prs):
    s = add_slide(prs, '자주 나오는 질문', '', 15)
    qa = [
        ('신호가 2배 강해지면 포지션도 2배가 되나요?',
         '아닙니다. 신호는 방향과 상대적 세기만 정하고, 크기는 변동성과 북 전체 '
         '위험 한도가 정합니다. 확신이 강해도 시장이 요동치면 포지션은 오히려 줄어듭니다.'),
        ('왜 영국·일본·호주는 신호만 보고 사지 않나요?',
         '성과가 나빠서가 아니라 매매 시간 때문입니다. 한국에서 그 시장들의 종가에 '
         '맞춰 주문하기 어렵습니다. 대신 신호 계산에는 계속 참여시켜 "한국·미국이 '
         '남들보다 나은가"를 판단하는 기준선 역할을 합니다.'),
        ('오늘 신호가 확 바뀌면 포지션도 확 바뀌나요?',
         '아닙니다. 새 목표의 20%만 반영합니다. 신호가 하루 튀어도 실제 주문은 '
         '조금씩만 나가고, 진짜 방향이 바뀌면 며칠에 걸쳐 따라갑니다.'),
        ('손실이 계속 나면 누가 멈추나요?',
         '북스톱이 자동으로 멈춥니다. 손실 4% 초과 시 절반, 8% 초과 시 전량 청산이며, '
         '회복하면 자동으로 다시 들어갑니다. 사람이 손으로 켜고 끄지 않습니다.'),
    ]
    y = Inches(1.30)
    for q, a in qa:
        _rect(s, Inches(0.45), y, Inches(12.4), Inches(1.10),
              fill=CARD_BG, line=CARD_LN)
        _rect(s, Inches(0.45), y, Inches(0.045), Inches(1.10), fill=ORANGE)
        tf = _textbox(s, Inches(0.70), y + Inches(0.11), Inches(12.0),
                      Inches(0.92))
        p = _para(tf, True)
        _run(p, 'Q.  ', size=12, bold=True, color=ORANGE)
        _run(p, q, size=12, bold=True, color=NAVY)
        p = _para(tf, False, space_before=5, line=1.25)
        _run(p, 'A.  ' + a, size=11, color=INK)
        y += Inches(1.32)


def s16_close(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    _rect(s, Inches(0.9), Inches(2.3), Inches(1.5), Inches(0.075), fill=ORANGE)
    tf = _textbox(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.0))
    p = _para(tf, True)
    _run(p, '기억할 것 3가지', size=32, bold=True, color=WHITE)

    items = [
        ('4명에게 묻고 평균 낸다',
         '추세·밸류·캐리·정책이 각각 25%씩. 어느 하나를 더 믿지 않습니다.'),
        ('신호는 방향, 리스크는 크기',
         '확신이 강해도 시장이 흔들리면 포지션은 작아집니다. 이 분리가 핵심입니다.'),
        ('천천히 움직이고, 알아서 멈춘다',
         '매일 20%씩만 조정하고, 손실이 커지면 북스톱이 자동으로 규모를 줄입니다.'),
    ]
    y = Inches(3.6)
    for i, (head, body) in enumerate(items):
        _rect(s, Inches(0.92), y, Inches(0.04), Inches(0.72), fill=ORANGE)
        tf2 = _textbox(s, Inches(1.20), y - Inches(0.03), Inches(11.0), Inches(0.8))
        p = _para(tf2, True)
        _run(p, f'{i + 1}.  {head}', size=16, bold=True, color=WHITE)
        p = _para(tf2, False, space_before=2)
        _run(p, body, size=11.5, color=RGBColor(0xBF, 0xD4, 0xEA))
        y += Inches(0.98)

    tf3 = _textbox(s, Inches(0.9), Inches(6.8), Inches(11.5), Inches(0.4))
    p = _para(tf3, True)
    _run(p, '더 자세한 내용: 「글로벌 금리 팩터 포트폴리오 — 운용 로직 설명서」',
         size=10, color=RGBColor(0x8F, 0xAD, 0xCC))


BUILDERS = [s1_cover, s2_signconv, s3_onesentence, s4_trend, s5_value, s6_carry,
            s7_policy, s8_compare, s9_jargon, s10_step1, s11_step2, s12_step3,
            s13_step4, s14_final, s15_faq, s16_close]


def build(out_path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    for fn in BUILDERS:
        fn(prs)
    prs.save(out_path)
    print(f"[저장] {out_path}  ({len(BUILDERS)} slides)")


if __name__ == '__main__':
    ap = argparse.ArgumentParser()
    ap.add_argument('-o', '--out',
                    default=f'금리팩터_쉬운설명_{date.today().isoformat()}.pptx')
    build(ap.parse_args().out)
