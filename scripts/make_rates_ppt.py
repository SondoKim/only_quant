# -*- coding: utf-8 -*-
"""금리 팩터 포트폴리오 운용 로직 설명 PPT 생성 (팀 운용 전환 브리핑용).

미래에셋 CI 색상(NAVY #043B72 / ORANGE #F58220 — total_dashboard/app.py 와 동일)
을 쓰는 12장 덱. 내용은 sleeve_engine.py / config/indicators.yaml 의 현행 운용
설정 기준이며, 숫자를 바꿀 일이 생기면 이 파일의 상수·표 데이터만 고치면 된다.

⚠ 2026-07-24 볼타겟 북 분리(금리 7.1% / FX 7.1%) 반영 — 금리 SR 1.15 기준.
   결합 모드 시절 수치(SR 1.20~1.21)를 다시 쓰지 말 것.

Usage: python scripts/make_rates_ppt.py [-o 출력경로.pptx]
"""
import argparse
from datetime import date

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ── 미래에셋 CI ──────────────────────────────────────────────────────────
NAVY    = RGBColor(0x04, 0x3B, 0x72)
ORANGE  = RGBColor(0xF5, 0x82, 0x20)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x1F, 0x2A, 0x37)   # 본문
MUTED   = RGBColor(0x6B, 0x72, 0x80)   # 보조 설명
CARD_BG = RGBColor(0xF6, 0xF8, 0xFA)
CARD_LN = RGBColor(0xDF, 0xE4, 0xEA)
RED     = RGBColor(0xC0, 0x39, 0x2B)   # 경고·금지
GREEN   = RGBColor(0x1E, 0x7A, 0x4C)   # 양호

FONT = '맑은 고딕'
W, H = Inches(13.333), Inches(7.5)


# ── 저수준 헬퍼 ──────────────────────────────────────────────────────────
def _run(p, text, size=11, bold=False, color=INK, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = FONT
    return r


def _para(tf, first, space_before=0, align=PP_ALIGN.LEFT, line=1.0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.alignment = align
    p.line_spacing = line
    return p


def _rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def _textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    return tf


# ── 슬라이드 골격 ────────────────────────────────────────────────────────
def add_slide(prs, title, subtitle=None, page=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    _rect(slide, 0, 0, W, Inches(0.92), fill=NAVY)
    _rect(slide, 0, Inches(0.92), W, Inches(0.055), fill=ORANGE)

    tf = _textbox(slide, Inches(0.45), Inches(0.10), Inches(11.0), Inches(0.72),
                  anchor=MSO_ANCHOR.MIDDLE)
    p = _para(tf, True)
    _run(p, title, size=21, bold=True, color=WHITE)
    if subtitle:
        p2 = _para(tf, False, space_before=1)
        _run(p2, subtitle, size=10.5, color=RGBColor(0xBF, 0xD4, 0xEA))

    if page is not None:
        tfp = _textbox(slide, W - Inches(0.95), H - Inches(0.42),
                       Inches(0.6), Inches(0.28))
        p = _para(tfp, True, align=PP_ALIGN.RIGHT)
        _run(p, str(page), size=9, color=MUTED)
    return slide


def card(slide, x, y, w, h, head, lines, accent=ORANGE, head_size=13.5,
         body_size=11, bg=CARD_BG):
    """제목 + 불릿 카드. lines = [str | (str, color) | ('•', str)]"""
    _rect(slide, x, y, w, h, fill=bg, line=CARD_LN)
    _rect(slide, x, y, Inches(0.045), h, fill=accent)

    tf = _textbox(slide, x + Inches(0.20), y + Inches(0.11),
                  w - Inches(0.33), h - Inches(0.20))
    first = True
    if head:
        p = _para(tf, first)
        first = False
        _run(p, head, size=head_size, bold=True, color=NAVY)
    for item in lines:
        color, txt, bold = INK, item, False
        if isinstance(item, tuple):
            if len(item) == 3:
                txt, color, bold = item
            else:
                txt, color = item
        p = _para(tf, first, space_before=9 if not first else 5, line=1.3)
        first = False
        _run(p, txt, size=body_size, color=color, bold=bold)
    return tf


def table(slide, x, y, w, headers, rows, col_w=None, fs=9.5, hfs=9.5,
          row_h=Inches(0.30), head_h=Inches(0.34), cell_colors=None):
    """헤더 네이비 / 본문 흰색 줄무늬 표. cell_colors[r][c] = RGBColor|None"""
    nr, nc = len(rows) + 1, len(headers)
    shape = slide.shapes.add_table(nr, nc, x, y, w, head_h + row_h * len(rows))
    tbl = shape.table
    tbl.first_row = True
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Inches(w.inches * cw / total)
    tbl.rows[0].height = head_h
    for i in range(len(rows)):
        tbl.rows[i + 1].height = row_h

    for c, htxt in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.06)
        tf = cell.text_frame
        tf.word_wrap = True
        p = _para(tf, True, align=PP_ALIGN.CENTER)
        _run(p, htxt, size=hfs, bold=True, color=WHITE)

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 0 else RGBColor(0xF4, 0xF7, 0xFA)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.06)
            tf = cell.text_frame
            tf.word_wrap = True
            col = INK
            if cell_colors and cell_colors.get((r, c)):
                col = cell_colors[(r, c)]
            p = _para(tf, True, align=PP_ALIGN.CENTER if c else PP_ALIGN.LEFT)
            _run(p, str(val), size=fs, color=col, bold=(c == 0))
    return tbl


def note(slide, y, text, color=MUTED, size=9, x=Inches(0.45), w=Inches(12.4)):
    tf = _textbox(slide, x, y, w, Inches(0.5))
    p = _para(tf, True, line=1.2)
    _run(p, text, size=size, color=color)
    return tf


# ═══════════════════════════════════════════════════════════════════════
# 슬라이드 정의
# ═══════════════════════════════════════════════════════════════════════
def s01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    _rect(slide, Inches(0.9), Inches(2.55), Inches(1.5), Inches(0.075), fill=ORANGE)

    tf = _textbox(slide, Inches(0.9), Inches(2.75), Inches(11.4), Inches(1.5))
    p = _para(tf, True)
    _run(p, '글로벌 금리 팩터 포트폴리오', size=38, bold=True, color=WHITE)
    p = _para(tf, False, space_before=6)
    _run(p, '운용 로직 설명서', size=25, color=ORANGE)

    tf2 = _textbox(slide, Inches(0.92), Inches(4.55), Inches(11.0), Inches(1.2))
    p = _para(tf2, True, line=1.35)
    _run(p, '4팩터(Trend · Value · Carry · Policy) 시그널로 '
            '한국·미국 국채선물을 운용하는 시스템 트레이딩 북',
         size=13, color=RGBColor(0xBF, 0xD4, 0xEA))
    p = _para(tf2, False, space_before=8, line=1.35)
    _run(p, '시그널 유니버스 8종 · 매매 4종 (미국 2Y/10Y · 한국 3Y/10Y)',
         size=13, color=RGBColor(0xBF, 0xD4, 0xEA))

    _rect(slide, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.02),
          fill=RGBColor(0x2A, 0x5A, 0x8C))
    tf3 = _textbox(slide, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.5))
    p = _para(tf3, True)
    _run(p, f'{date.today().isoformat()}   |   검증 기간 2016-01 ~ 2026-07 '
            f'(비용 차감 후)   |   내부 운용 자료',
         size=10, color=RGBColor(0x8F, 0xAD, 0xCC))


def s02_summary(prs):
    s = add_slide(prs, '한 장 요약 — 이 북은 무엇을 하는가',
                  '팩터가 방향을 정하고, 리스크 모델이 크기를 정한다', 2)
    y = Inches(1.28)
    card(s, Inches(0.45), y, Inches(4.0), Inches(2.55), '① 무엇을 보는가',
         ['글로벌 국채선물 8종의 4개 팩터를 매일 z-score로 측정',
          ('Trend · Value · Carry · Policy (가중 각 1.0)', NAVY, True),
          '금리(yield) 데이터는 시그널 전용 — 매매하지 않음'])
    card(s, Inches(4.65), y, Inches(4.0), Inches(2.55), '② 무엇을 매매하는가',
         [('미국 2Y · 미국 10Y · 한국 3Y · 한국 10Y', NAVY, True),
          '영국·일본·호주는 시그널에만 참여(집행시간 제약)',
          '주문은 t 종가 기준 산출 → t+1 집행'])
    card(s, Inches(8.85), y, Inches(4.0), Inches(2.55), '③ 크기는 어떻게 정하는가',
         ['팩터는 방향만 결정 — 신호가 2배여도 포지션 2배 아님',
          ('인버스볼 → 볼타겟 → 스무딩 → 북스톱', NAVY, True),
          '금리 북 연 7.1% 변동성 목표 (FX와 완전 분리)'])

    y2 = Inches(4.10)
    tf = _textbox(s, Inches(0.45), y2, Inches(12.4), Inches(0.35))
    p = _para(tf, True)
    _run(p, '검증 성과 (2016-01 ~ 2026-07, 거래비용 차감 후)', size=14, bold=True, color=NAVY)

    table(s, Inches(0.45), Inches(4.55), Inches(12.4),
          ['구분', 'Sharpe', '연 수익률', '변동성', 'MaxDD', 'T+2 지연 집행', '연 회전율'],
          [['금리 북', '1.15', '3.0%', '2.6%', '-4.7%', '1.08', '17x'],
           ['FX 북 (참고)', '0.34', '2.2%', '7.2%', '-20.3%', '—', '—'],
           ['합산', '0.69', '5.3%', '7.9%', '-21.9%', '—', '—']],
          col_w=[2.2, 1, 1.2, 1, 1, 1.6, 1.1], fs=11,
          row_h=Inches(0.42), head_h=Inches(0.40),
          cell_colors={(0, 1): GREEN, (0, 4): GREEN, (0, 5): GREEN})

    note(s, Inches(6.35),
         '▪ 금리 북 변동성 2.6%는 목표 7.1%에 노출스케일 0.5와 포지션 스무딩·북스톱이 함께 걸린 결과입니다 (실현 < 목표는 정상).\n'
         '▪ T+2(하루 늦게 집행)에서도 Sharpe 1.08 유지 — 이 성과가 집행 타이밍 아티팩트가 아니라는 핵심 증거입니다.')


def s03_universe(prs):
    s = add_slide(prs, '유니버스 — 8종을 보고 4종을 매매한다',
                  '시그널 범위와 매매 범위를 일부러 다르게 둔 구조', 3)
    table(s, Inches(0.45), Inches(1.30), Inches(6.15),
          ['자산', '역할'],
          [['미국 2Y (TU1)', '매매'], ['미국 10Y (TY1)', '매매'],
           ['한국 3Y (KE1)', '매매'], ['한국 10Y (KAA1)', '매매'],
           ['영국 10Y (G 1)', '시그널 전용'], ['일본 10Y (JB1)', '시그널 전용'],
           ['호주 3Y (YM1)', '시그널 전용'], ['호주 10Y (XM1)', '시그널 전용']],
          col_w=[2.0, 1.4], fs=11, row_h=Inches(0.345), head_h=Inches(0.38),
          cell_colors={(i, 1): (GREEN if i < 4 else MUTED) for i in range(8)})

    card(s, Inches(6.95), Inches(1.30), Inches(5.9), Inches(1.95),
         '왜 8종을 보면서 4종만 매매하나',
         ['매매 제외는 성과가 아니라 집행 가능성 때문 — 영국·일본·호주는 '
          '한국에서 매매 시간이 애매합니다.',
          ('시그널 유니버스를 4종으로 좁히면 오히려 악화 '
           '(Sharpe 1.20→1.11, MaxDD -5.5%→-7.3%)', NAVY, True)])

    card(s, Inches(6.95), Inches(3.45), Inches(5.9), Inches(2.85),
         '"반쪽 헤지" — 반드시 이해해야 할 구조',
         ['Value·Carry는 8종 평균 대비 상대값인데, 반대편 다리(영·일·호)는 '
          '주문하지 않습니다.',
          '→ 상대가치 신호가 완전히 상쇄되지 않고 순방향 노출로 일부 남습니다.',
          ('의도된 설계입니다. 넓은 횡단면이 주는 신호 품질이 '
           '반쪽 헤지의 손해보다 큽니다.', ORANGE, True)])

    card(s, Inches(0.45), Inches(4.55), Inches(6.15), Inches(1.75),
         '⚠ 구조적 취약점 — 실효 독립 베팅 ≈ 1.86개',
         ['ρ(미국2Y, 미국10Y) = 0.78 · ρ(한국3Y, 한국10Y) = 0.87',
          '→ 사실상 "미국 듀레이션 + 한국 듀레이션" 2개 베팅',
          ('한·미 금리가 동시에 역행하는 글로벌 쇼크에는 상쇄 수단이 없습니다.',
           RED, True)], accent=RED)

    note(s, Inches(6.48),
         '※ 유로존 4종(독일 2Y/10Y·프랑스·이탈리아)은 A/B 검증 후 제외했습니다 — 포함 시 Sharpe 0.88, 제외 시 1.20. '
         '앞으로 유니버스를 더 깎지 말 것 (12→8로 이미 선택편의가 누적).')


def s04_factor_a(prs):
    s = add_slide(prs, '4팩터 (1/2) — Trend · Value',
                  '방향성 팩터와 상대가치 팩터를 의도적으로 섞는다', 4)
    card(s, Inches(0.45), Inches(1.30), Inches(6.1), Inches(3.05),
         '① Trend (추세) — 가중 1.0 · 방향성',
         [('계산: z(126일 수익률)과 z(252일 수익률)의 평균, ±3 클립', NAVY, True),
          '논리: 금리 사이클은 수개월 단위 관성이 있다 (TSMOM).',
          '63일 같은 빠른 구간은 금리에서 휩소가 심해 의도적으로 배제.',
          ('레짐 게이트: 120일 Hurst ≤ 0.5(횡보장)이면 추세 신호를 0으로 차단',
           ORANGE, True),
          '횡단면 중립화 안 함 → 8개국이 함께 오르면 함께 롱. 이 방향성 '
          '노출이 이 북 알파의 주 원천.'])
    card(s, Inches(6.75), Inches(1.30), Inches(6.1), Inches(3.05),
         '② Value (밸류) — 가중 1.0 · 완전 시장중립',
         [('계산: −z(log 가격, 504일), ±3 클립 → 8종 횡단면 평균 차감', NAVY, True),
          '논리: 2년 평균 대비 비싸진 채권은 판다 (장기 평균회귀).',
          '완전 중립화의 의미 — 오늘 미국 2Y는 절대값으로는 "싸다(+1.73)"고 '
          '나오지만, 8개국이 더 싸서 상대적으로는 비싼 편(−0.26)입니다.',
          ('Value는 "채권이 싼가"가 아니라 "어느 나라가 남들보다 싼가"만 '
           '말합니다.', ORANGE, True)])

    card(s, Inches(0.45), Inches(4.55), Inches(12.4), Inches(1.75),
         '방향성(xs=0)과 상대가치(xs=1)를 나누는 이유',
         ['Trend·Policy는 횡단면 중립화를 하지 않아 "글로벌 듀레이션 방향" 베팅을 그대로 남깁니다. '
          'Value·Carry는 완전 중립화해 "국가 간 상대가치"만 남깁니다.',
          ('2026-07-23 A/B: Trend·Policy에 부분 중립화(0.25/0.5/1.0)를 얹자 '
           'Sharpe가 1.20 → 1.13~0.72로 단조 악화 — 전 격자 기각. '
           '이 북의 알파는 방향성 CTA 프리미엄입니다.', NAVY, True)])

    note(s, Inches(6.48),
         '※ 파라미터는 모두 사전등록 값이며 같은 표본에서 재튜닝하지 않습니다. 변경하려면 새 표본으로 A/B 후 config에 근거를 남기는 것이 이 북의 규약입니다.')


def s05_factor_b(prs):
    s = add_slide(prs, '4팩터 (2/2) — Carry · Policy',
                  '금리(yield) 데이터를 쓰는 두 팩터', 5)
    card(s, Inches(0.45), Inches(1.30), Inches(6.1), Inches(3.05),
         '③ Carry (캐리) — 가중 1.0 · 완전 시장중립',
         [('계산: 자국 만기금리 − (자국 단기금리 − 미국 2Y), 252일 z', NAVY, True),
          '= 달러 투자자가 환헤지 비용까지 차감하고 실제로 버는 캐리 (CIP 근사).',
          '자국 단기금리가 높아 헤지비용이 다 먹히는 시장은 "공짜 캐리"로 '
          '잡히지 않습니다.',
          ('→ 8종 횡단면 평균 차감 후 "고캐리국 롱 / 저캐리국 숏"', ORANGE, True)])
    card(s, Inches(6.75), Inches(1.30), Inches(6.1), Inches(3.05),
         '④ Policy (정책 모멘텀) — 가중 1.0 · 방향성',
         [('계산: −(자국 2Y/3Y 금리 − 126일 전 값), 252일 z', NAVY, True),
          '논리: 단기금리가 6개월간 하락 = 완화 사이클 → 듀레이션 롱.',
          '상승 = 긴축 사이클 → 숏.',
          ('전략 공장에서 가장 강했던 alpha4 policy_momentum의 연속형 버전.',
           ORANGE, True),
          '중앙은행 사이클을 가격이 아니라 금리로 직접 읽는 유일한 팩터.'])

    card(s, Inches(0.45), Inches(4.55), Inches(12.4), Inches(1.90),
         '⚠ 팀이 알아야 할 구조적 특이점 (실측 검증)',
         [('Carry는 미국 2Y와 한국 3Y를 구분하지 못합니다 — 두 신호의 상관 1.000, '
           '최대 차이 9.5e-14(완전 동일).', RED, True),
          '만기가 펀딩 프록시와 같은 단기물은 헤지캐리가 수식상 미국 2Y로 수렴하기 때문입니다 '
          '(이론적으로는 옳지만, 이 팩터의 실효 정보량은 4자산이 아니라 3자산).',
          ('Policy는 나라만 구분하고 만기는 구분하지 않습니다 — 미국 2Y=미국 10Y, '
           '한국 3Y=한국 10Y (차이 정확히 0.0).', RED, True)], accent=RED)

    note(s, Inches(6.62),
         '※ 커브 팩터(10Y−2Y 기울기)는 가중 0 = 비활성입니다. 2026-06-11 및 2026-07-23 두 차례 A/B에서 모두 Sharpe를 깎아 기각됐습니다.')


def s06_pipeline(prs):
    s = add_slide(prs, '팩터에서 주문까지 — 8단계 파이프라인',
                  '모든 단계는 causal (t 시점 포지션은 t까지의 데이터만 사용)', 6)
    steps = [
        ('1', '팩터 결합', '4팩터 단순 평균 (가중 전부 1.0)\n→ EWMA span 5로 잡음 제거',
         '동일 가중 = 특정 팩터 과최적화 방지'),
        ('2', '매매 제외 적용', '영·일·호 포지션만 0으로\n신호 계산에서 빼는 게 아님',
         '이들이 횡단면 기준선을 계속 만듦'),
        ('3', '인버스볼 사이징', '0.10 / EWMA(halflife 33일) 변동성\n저변동 자산이 더 많은 계약',
         '개별 포지션 상한 ±3.0'),
        ('4', '포트폴리오 볼타겟', '과거 63일 수익률에 오늘 비중 적용\n금리 북 단독 연 7.1% 목표',
         'FX 북과 완전 분리 (2026-07-24)'),
        ('5', '포지션 스무딩 0.8', '오늘 신규 목표의 20%만 반영\n회전율 134x→73x, Sharpe 0.40→0.51',
         '급격한 신호 변화를 며칠에 걸쳐 반영'),
        ('6', '북스톱', '섀도우 손실 4% 초과 → 절반\n8% 초과 → 플랫',
         '섀도우 기준이라 회복 시 자동 재진입'),
        ('7', '노출 스케일 0.5', '운용 북 델타 규모 캘리브레이션\n최종 포지션에만 곱함',
         '시그널·스톱 동작은 그대로, 규모만 축소'),
        ('8', '주문', 't 종가 기준 산출 → t+1 집행\n연 회전 17x (일 주문 소량)',
         'T+2로 밀려도 Sharpe 1.08 유지'),
    ]
    x0, y0 = Inches(0.45), Inches(1.32)
    cw, ch, gap = Inches(3.02), Inches(2.35), Inches(0.11)
    for i, (num, head, body, effect) in enumerate(steps):
        col, row = i % 4, i // 4
        x = x0 + (cw + gap) * col
        y = y0 + (ch + Inches(0.34)) * row
        _rect(s, x, y, cw, ch, fill=CARD_BG, line=CARD_LN)
        _rect(s, x, y, cw, Inches(0.045), fill=ORANGE if row == 0 else NAVY)
        badge = _rect(s, x + Inches(0.13), y + Inches(0.17), Inches(0.30),
                      Inches(0.30), fill=NAVY, shape=MSO_SHAPE.OVAL)
        btf = badge.text_frame
        btf.margin_left = btf.margin_right = 0
        btf.margin_top = btf.margin_bottom = 0
        bp = _para(btf, True, align=PP_ALIGN.CENTER)
        _run(bp, num, size=10, bold=True, color=WHITE)

        tf = _textbox(s, x + Inches(0.52), y + Inches(0.14), cw - Inches(0.65),
                      Inches(0.34))
        p = _para(tf, True)
        _run(p, head, size=11.5, bold=True, color=NAVY)
        tf2 = _textbox(s, x + Inches(0.16), y + Inches(0.56), cw - Inches(0.30),
                       ch - Inches(0.68), anchor=MSO_ANCHOR.MIDDLE)
        p = _para(tf2, True, line=1.3)
        _run(p, body, size=9.5, color=INK)
        p = _para(tf2, False, space_before=11, line=1.3)
        _run(p, effect, size=9, color=ORANGE, bold=True)

    note(s, Inches(6.55),
         '▪ 4단계(볼타겟)는 2026-07-24부터 금리 북과 FX 북을 완전 분리했습니다 — FX 변동성이 튀어도 금리 포지션이 흔들리지 않습니다.\n'
         '▪ 5단계 스무딩 때문에 팩터 신호가 크게 흔들려도 실제 주문은 잘게 나옵니다 (예: 미국 2Y −0.79 → −0.77).')


def s07_risk(prs):
    s = add_slide(prs, '리스크 관리 — 3중 구조', '크기 / 총량 / 손실 각각에 독립 장치', 7)
    card(s, Inches(0.45), Inches(1.30), Inches(4.0), Inches(2.85),
         '① 인버스볼 — 자산별 등리스크',
         [('자산당 연 10% 리스크 기준으로 계약수 결정', NAVY, True),
          'EWMA(halflife 33일) 실현변동성의 역수로 사이징.',
          '변동성이 낮은 미국 2Y는 더 많은 계약, 변동성 높은 자산은 더 적게.',
          '개별 포지션 상한 ±3.0'])
    card(s, Inches(4.65), Inches(1.30), Inches(4.0), Inches(2.85),
         '② 포트폴리오 볼타겟 — 총량',
         [('금리 북 단독 연 7.1% 변동성 목표', NAVY, True),
          '오늘 비중을 과거 63일 수익률에 적용 → look-ahead 없음.',
          ('2026-07-24부터 FX 북과 완전 분리 — 검증: FX 타겟을 30%로 바꿔도 '
           '금리 포지션 변화 정확히 0.0', ORANGE, True),
          '초기 63거래일은 강제 플랫(워밍업)'])
    card(s, Inches(8.85), Inches(1.30), Inches(4.0), Inches(2.85),
         '③ 북스톱 — 손실 제어',
         [('섀도우 손실 4% 초과 → 포지션 절반\n8% 초과 → 플랫', NAVY, True),
          '"스톱이 없었다면 겪었을" 가상 손실을 봅니다.',
          ('그래서 플랫 상태에서도 회복하면 자동 재진입 — 수동으로 켜지 마십시오.',
           RED, True),
          '고점 기준 = 252일 롤링 (전기간 아님)'])

    card(s, Inches(0.45), Inches(4.35), Inches(12.4), Inches(1.85),
         '북스톱이 252일 롤링 고점을 쓰는 이유 (2026-07-22 감사)',
         ['원설계는 전기간 고점 대비로 손실을 재서, 섀도우 북이 옛 고점(2022-09-28)을 못 넘자 '
          '드로다운이 임계 위에 영구 고착됐습니다.',
          ('결과: 2023-07 이후 계속 플랫 — 손실 제어 장치가 사실상 킬스위치로 작동 '
           '(2016년 이후 42% 일수 플랫).', RED, True),
          '이제 오래된 고점이 시간이 지나면 굴러떨어지므로 재진입 경로가 항상 열려 있습니다.'],
         accent=RED)

    note(s, Inches(6.42),
         '※ 임계값 4%/8%은 사전등록 프로토타입 값입니다. 스윕에서 3%/6%이 더 높게 나왔지만 같은 표본에서 재튜닝하지 않았습니다 — 메커니즘 버그만 고치고 파라미터는 건드리지 않는 것이 이 북의 원칙입니다.')


def s08_example(prs):
    s = add_slide(prs, '오늘의 북 — 실제 숫자로 읽기',
                  '2026-07-23 종가 기준 · 대시보드가 자동 생성하는 해석과 동일', 8)
    table(s, Inches(0.45), Inches(1.28), Inches(7.5),
          ['구분', '미국 2Y', '미국 10Y', '한국 3Y', '한국 10Y'],
          [['Trend (방향성)', '-2.19', '-1.96', '-1.14', '-1.25'],
           ['Value (중립화)', '-0.26', '+0.07', '+0.34', '+0.51'],
           ['Carry (중립화)', '-0.02', '-0.00', '-0.02', '+0.25'],
           ['Policy (방향성)', '-1.96', '-1.96', '-0.91', '-0.91'],
           ['컨빅션 (4팩터 평균)', '-1.11', '-0.96', '-0.43', '-0.35'],
           ['최종 목표 포지션', '-0.77', '-0.50', '-0.28', '-0.07']],
          col_w=[2.3, 1.2, 1.2, 1.2, 1.2], fs=11, row_h=Inches(0.42),
          head_h=Inches(0.40), cell_colors={(5, c): NAVY for c in range(5)})

    card(s, Inches(8.25), Inches(1.28), Inches(4.6), Inches(2.92),
         '🧭 자동 생성 해석',
         [('4종 전부 듀레이션 숏 (순 −1.61 / 그로스 1.61 = 100% 방향성 베팅)',
           NAVY, True),
          '근거: Trend(전 만기 하락추세)와 Policy(미국 2Y가 6개월간 상승 = '
          '긴축 신호 −1.96)가 둘 다 숏을 가리킵니다.',
          'Value·Carry는 한국 쪽에 소폭 롱 기울기를 주지만(한국 10Y +0.51/+0.25) '
          '방향성 두 팩터를 못 이깁니다.',
          ('결과: 미국 숏이 크고(−0.77/−0.50) 한국 숏은 작습니다(−0.28/−0.07).',
           ORANGE, True)])

    card(s, Inches(0.45), Inches(4.45), Inches(12.4), Inches(1.65),
         '읽을 때 주의 — 컨빅션 ≠ 포지션 크기',
         ['컨빅션(−1.11)과 최종 포지션(−0.77)은 비례하지 않습니다. '
          '사이에 인버스볼 → 볼타겟 → 스무딩 → 북스톱 → 노출스케일이 들어갑니다.',
          ('Value·Carry 값은 8종 평균 대비 상대값입니다 — 매매하지 않는 영·일·호가 '
           '기준선을 만들어 한·미 포지션을 실제로 좌우합니다.', ORANGE, True)])

    note(s, Inches(6.30),
         '※ 이 해석 문구는 대시보드가 슬리브 z값에서 코드로 자동 생성합니다 (AI 생성 아님 — 같은 입력이면 항상 같은 문장, 20회 실행 해시 동일 검증). '
         '「매크로 전략 공장 → 자산별 트레이딩 시그널 → 금리 자산 시그널」 표 아래에서 매일 확인할 수 있습니다.')


def s09_evidence(prs):
    s = add_slide(prs, '이 성과를 믿어도 되는 이유', '아티팩트를 걸러낸 검증 프로토콜', 9)
    table(s, Inches(0.45), Inches(1.28), Inches(7.4),
          ['검증 항목', '결과', '판정'],
          [['기본 (T+1 집행, 비용 차감)', 'Sharpe 1.15', '기준'],
           ['T+2 — 하루 늦게 집행', 'Sharpe 1.08', '통과'],
           ['시간대 정직 (조기마감 T+2/미국 T+1)', 'Sharpe 1.14', '통과'],
           ['전반기 / 후반기 분할', '1.27 / 1.02', '통과'],
           ['장기 표본 (2012년~)', 'Sharpe 0.77', '통과'],
           ['최대 낙폭', '-4.7%', '통과']],
          col_w=[3.2, 1.3, 0.9], fs=11, row_h=Inches(0.42), head_h=Inches(0.40),
          cell_colors={(i, 2): GREEN for i in range(1, 6)})

    card(s, Inches(8.15), Inches(1.28), Inches(4.7), Inches(2.92),
         '왜 T+2 테스트가 결정적인가',
         ['한국·일본 종가는 시카고 종가보다 약 14시간 이릅니다. '
          '아시아 가격은 미국 움직임을 하루 늦게 반영합니다.',
          ('과거 이 북에 있던 "단기 평균회귀" 서브북은 바로 이 시차를 수확하고 '
           '있었습니다 — 하루만 늦추면 Sharpe 1.01 → −0.14로 소멸.', RED, True),
          '미국 종가를 보고 한국 종가에 체결하는 것은 물리적으로 불가능하므로 '
          '2026-07-22 영구 비활성했습니다.',
          ('현재 북은 T+2에서도 1.08 — 시차 아티팩트가 아닙니다.', GREEN, True)])

    card(s, Inches(0.45), Inches(4.45), Inches(12.4), Inches(1.95),
         '검증 규약 — 팀이 함께 지켜야 할 원칙',
         ['① 파라미터는 사전등록 후 A/B — 같은 표본에서 결과를 보고 재튜닝하지 않습니다.',
          '② 부분 채택 금지 — "미국 커브만 켜자" 같은 사후 선택은 같은 표본 재튜닝입니다.',
          '③ 기각된 실험도 config 주석과 기록에 남깁니다 (2026-07-23 하루에만 커브 트레이드·블록 예산·횡단면 중립화 3건 기각).',
          ('④ 비용 가정 0.5bp는 이 북 성과의 전제입니다 — 실제 슬리피지가 이를 넘으면 성과는 재검증 대상입니다.', ORANGE, True)])


def s10_risks(prs):
    s = add_slide(prs, '운용 리스크 — 조용히 실패할 수 있는 경로',
                  '에러 없이 잘못된 신호가 나갈 수 있는 지점들', 10)
    card(s, Inches(0.45), Inches(1.28), Inches(6.1), Inches(2.55),
         '⚠ 1순위 — 금리 데이터 스테일',
         [('yield 패널은 결측 시 자동으로 전일값을 재사용(ffill)합니다.', RED, True),
          '블룸버그 피드가 끊기면 에러 없이 조용히 Carry·Policy가 고착됩니다.',
          ('현재 자동 경보가 없습니다 — 매일 yield 패널 최종 갱신일자를 '
           '확인하는 절차가 필요합니다.', RED, True)], accent=RED)
    card(s, Inches(6.75), Inches(1.28), Inches(6.1), Inches(2.55),
         '⚠ 2순위 — 데이터 리비전 / 워밍업',
         ['매일 전체 이력을 처음부터 재계산합니다 (상태 비저장).',
          '블룸버그가 과거 데이터를 수정하면 과거 포지션까지 바뀝니다 — '
          '재현 검증하려면 날짜별 원본 스냅샷 보관이 필요합니다.',
          ('Value 504일 + 볼 63일 → 최소 2.5년 이력 필요. '
           '데이터 시작일을 임의로 자르지 마십시오.', RED, True)], accent=RED)

    card(s, Inches(0.45), Inches(4.05), Inches(12.4), Inches(2.55),
         '⛔ 절대 켜지 말 것 — 모두 검증 후 기각된 항목',
         [('reversion (단기 평균회귀 서브북)  —  수익이 알파가 아니라 비동시 종가 '
           '아티팩트. 켜면 이후 모든 A/B가 오염됩니다.', RED, True),
          ('curve_trades (합성 커브 트레이드)  —  2026-07-23 A/B 기각 '
           '(Sharpe 1.20 → 0.70)', RED, True),
          ('risk_blocks (블록 리스크 예산)  —  2026-07-23 A/B 기각 '
           '(Sharpe 1.20 → 1.01, MaxDD 악화)', RED, True),
          ('xs_neutralize의 trend/policy를 0 초과로  —  2026-07-23 A/B 기각 '
           '(전 격자 단조 악화, 최악 0.72)', RED, True),
          ('유니버스 추가 축소  —  12→8로 이미 선택편의가 누적돼 있습니다.',
           RED, True)], accent=RED)


def s11_checklist(prs):
    s = add_slide(prs, '일일 운영 체크리스트', '매일 아침 5분, 순서대로', 11)
    table(s, Inches(0.45), Inches(1.28), Inches(12.4),
          ['#', '점검 항목', '무엇을 보는가', '이상 시 조치'],
          [['1', '데이터 신선도', '가격 18종 · 금리 15종의 최종일자가 오늘인가',
            '스테일이면 주문 보류 후 피드 확인'],
           ['2', '포지션 변화 폭', '목표 포지션 vs 전일 — 통상 자산당 ±0.05 수준',
            '급변 시 데이터 이상 여부 우선 확인'],
           ['3', '북스톱 상태', '대시보드 해석의 "북스톱" 줄이 정상 / 감축 / 발동 중 무엇인가',
            '감축·발동이면 팀 공지 (수동 복구 금지)'],
           ['4', '시그널 해석 읽기', '오늘 방향의 근거가 납득되는가',
            '납득 안 되면 팩터 시계열 차트 대조'],
           ['5', '집행 슬리피지', '실제 체결이 0.5bp 가정 안에 드는가',
            '초과 지속 시 성과 전제 재검토'],
           ['6', '집행 타이밍', 't 종가 기준 신호 → t+1 집행',
            '하루 밀려도 Sharpe 1.08 — 서두르지 말 것']],
          col_w=[0.4, 2.0, 4.6, 3.4], fs=10.5, row_h=Inches(0.50), head_h=Inches(0.40))

    card(s, Inches(0.45), Inches(4.92), Inches(6.1), Inches(1.72),
         '재량 개입에 대한 원칙',
         [('북스톱이 발동했을 때 수동으로 포지션을 되돌리지 마십시오.', RED, True),
          '섀도우 북이 회복하면 자동 재진입하도록 설계돼 있어, 손으로 켜면 '
          '로직이 깨지고 이후 성과 귀속이 불가능해집니다.'], accent=RED)
    card(s, Inches(6.75), Inches(4.92), Inches(6.1), Inches(1.72),
         '다음 개선 후보 (새 데이터 필요)',
         ['기존 가격·금리 데이터 안에서의 구조 개선 여지는 사실상 소진됐습니다.',
          ('① 인플레이션(CPI) 팩터 — 월간 데이터 파이프라인 추가\n'
           '② 미국 오버나이트 → 한국 개장 집행 연구 — KTB 개장가 데이터 확보',
           NAVY, True)])


def s12_close(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    _rect(slide, Inches(0.9), Inches(2.3), Inches(1.5), Inches(0.075), fill=ORANGE)

    tf = _textbox(slide, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.0))
    p = _para(tf, True)
    _run(p, '핵심 3가지', size=32, bold=True, color=WHITE)

    items = [
        ('팩터는 방향, 리스크 모델은 크기',
         '신호가 2배 강해져도 포지션은 2배가 되지 않습니다. 이 분리가 설계의 핵심입니다.'),
        ('이 북의 알파는 방향성 CTA 프리미엄',
         '상대가치화 시도(커브·블록예산·횡단면 중립화)는 모두 검증에서 기각됐습니다.'),
        ('가장 큰 실패 경로는 시장이 아니라 데이터',
         '금리 피드가 끊겨도 시스템은 에러 없이 전일값으로 신호를 냅니다. 매일 확인이 필요합니다.'),
    ]
    y = Inches(3.65)
    for i, (head, body) in enumerate(items):
        _rect(slide, Inches(0.92), y, Inches(0.04), Inches(0.72), fill=ORANGE)
        tf2 = _textbox(slide, Inches(1.20), y - Inches(0.03), Inches(11.0), Inches(0.8))
        p = _para(tf2, True)
        _run(p, f'{i + 1}.  {head}', size=15, bold=True, color=WHITE)
        p = _para(tf2, False, space_before=2)
        _run(p, body, size=11, color=RGBColor(0xBF, 0xD4, 0xEA))
        y += Inches(0.95)

    tf3 = _textbox(slide, Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.4))
    p = _para(tf3, True)
    _run(p, '상세 로직·검증 기록: config/indicators.yaml 주석 · '
            'src/sleeves/sleeve_engine.py · scripts/test_*.py',
         size=9.5, color=RGBColor(0x8F, 0xAD, 0xCC))


BUILDERS = [s01_cover, s02_summary, s03_universe, s04_factor_a, s05_factor_b,
            s06_pipeline, s07_risk, s08_example, s09_evidence, s10_risks,
            s11_checklist, s12_close]


def build(out_path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    for fn in BUILDERS:
        fn(prs)
    prs.save(out_path)
    print(f"[저장] {out_path}  ({len(BUILDERS)} slides)")


if __name__ == '__main__':
    ap = argparse.ArgumentParser()
    ap.add_argument('-o', '--out',
                    default=f'금리팩터_운용로직_{date.today().isoformat()}.pptx')
    build(ap.parse_args().out)
